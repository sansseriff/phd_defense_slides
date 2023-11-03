import collections
import collections.abc
import pptx
import win32com.client as client
import os
import shutil
import vss
from pyuac import main_requires_admin
import time
import subprocess
import argparse



def prev_slide(number):
    if number > 0:
        return number - 1
    else:
        return 0


def next_slide(number, ln):
    if number < ln - 1:
        return number + 1
    else:
        return ln - 1
    

# old

# <script>
#   import {{ onMount }} from 'svelte';
#   import {{ writable }} from 'svelte/store';
#   import {{ goto }} from '$app/navigation';
#   // <img src="/slides_png/slide_{i}.png" alt="slide_{i}" width="1000" height="500">

#   const width = writable(0);
#   const height = writable(0);
#   onMount(() => {{
#     width.set(window.innerWidth);
#     height.set(window.innerHeight);
#     document.addEventListener('keydown', async function(event) {{
#       //Check if the pressed key is the left arrow key
#       if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
#         // Navigate to the desired URL when the left arrow key is pressed
#         //window.location.href = '/slide_{prev_slide(i)}'; // Replace with your desired URL
#         await goto('/slide_{prev_slide(i)}');
#       }}
#       //Check if the pressed key is the right arrow key
#       else if (event.key === 'ArrowRight' || event.key === 'PageDown') {{
#         // Navigate to the desired URL when the right arrow key is pressed
#         //window.location.href = '/slide_{next_slide(i, ln)}'; // Replace with your desired URL
#         await goto('/slide_{next_slide(i, ln)}');
#       }}
#     }});
#   }});
# </script>

def create_svelte_route(i, ln, routes_directory):
    # I should add a checker here that looks in the route directory
    # if its there and does not overwrite it if theres some keyfile 
    # indicating dynamic content

    folder_name = f"slide_{i}"
    folder_path = os.path.join(routes_directory, folder_name)
    os.makedirs(folder_path, exist_ok=True)

    # Create a text file in the folder
    svelte_path = os.path.join(folder_path, "+page.svelte")
    with open(svelte_path, "w") as f:
        f.write(
            f"""
<script lang='ts'>
  import {{ onMount, onDestroy }} from 'svelte';
  import {{ writable }} from 'svelte/store';
  import {{ goto }} from '$app/navigation';

  const width = writable(0);
  const height = writable(0);

  let keydownHandler: (event: KeyboardEvent) => void;

  onMount(() => {{
    width.set(window.innerWidth);
    height.set(window.innerHeight);

    keydownHandler = async function(event) {{
      // Define an async function
      const navigate = async (url: string, imgSrc: string) => {{
        const img = new Image();
        img.src = imgSrc;
        img.onload = async () => {{
          await goto(url);
        }};
      }}
      //Check if the pressed key is the left arrow key
      if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
        // Navigate to the desired URL when the left arrow key is pressed
        navigate('/slide_{prev_slide(i)}', '/slides_png/slide_{prev_slide(i)}.png');
      }}
      //Check if the pressed key is the right arrow key
      else if (event.key === 'ArrowRight' || event.key === 'PageDown') {{
        // Navigate to the desired URL when the right arrow key is pressed
        navigate('/slide_{next_slide(i, ln)}', '/slides_png/slide_{next_slide(i, ln)}.png');
      }}
    }};

    document.addEventListener('keydown', keydownHandler);
  }});

  onDestroy(() => {{
    if (typeof window !== 'undefined') {{
      document.removeEventListener('keydown', keydownHandler);
    }}
  }});
</script>


<img src="/slides_png/slide_{i}.png" alt="slide_{i}" width="{{$width}}">
              
"""
        )

    js_path = os.path.join(folder_path, "+page.js")
    with open(js_path, "w") as f:
        f.write(
            """
// since there's no dynamic data here, we can prerender
// it so that it gets served as a static asset in production
export const prerender = true;
"""
        )
        # end

# @main_requires_admin(return_output=False)
# def main():
#     file_export = "./powerpoint/defense_export.pptx"
#     file = "./powerpoint/defense.pptx"
    

#     # shutil.copyfile(file, file_export)

#     full_path = os.path.abspath(file)


#     local_drives = set()
#     local_drives.add('C')
#     # Initialize the Shadow Copies
#     sc = vss.ShadowCopy(local_drives)
#     # An open and locked file we want to read
#     shadow_path = sc.shadow_path(full_path)

#     print(shadow_path)

#     shutil.copy(shadow_path, file_export)

#     print('testing')

file_export = "./powerpoint/defense_export.pptx"
export_path = ".\\web\\defense\\static\\slides_png\\"
routes_directory = "./web/defense/src/routes/"
full_export_path = os.path.abspath(export_path)

full_file_path = os.path.abspath(file_export)

def run_copy():
    subprocess.run(["python", "copier.py"])

def finish_up(run_copy_flag=True):
    # if run_copy_flag:
    run_copy()
        
    


    # works
    # ppt=Presentation("/Users/Andrew/defense.pptx")
    





    # Define the base directory where you want to create the folders
    

    # Define the number of folders you want to create

    ppt = pptx.Presentation(file_export)
    notes = []
    # https://stackoverflow.com/questions/61815883/python-pptx-export-img-png-jpeg

    Application = client.Dispatch("PowerPoint.Application")
    Presentation = Application.Presentations.Open(full_file_path)


    # Loop through the range of folder numbers and create each folder
    for i, (note_slide, img_slide) in enumerate(zip(ppt.slides, Presentation.Slides)):
        textNote = note_slide.notes_slide.notes_text_frame.text
        notes.append((i, textNote))

        if "dyno" not in textNote:
            img_slide.Export(f"{full_export_path}\\slide_{i}.png", "PNG")

            create_svelte_route(i, len(ppt.slides), routes_directory)

        else:
            print("found dynamic slide")

    with open(f"{full_export_path}\\notes.json", "w") as f:
        f.write(str(notes))

    

    

    # Application.Quit() # this closes all instances of powerpoint.
    Presentation.Close() # this just closes the auto-opened instance
    Presentation = None
    Application = None


def fix_svelte_routes():
    with open(f"{full_export_path}\\notes.json", "r") as f:
        notes = eval(f.read())

    # print(notes)

    for i, note in notes:
        if "dyno" not in note:
            create_svelte_route(i, len(notes), routes_directory)
        else:
            print("found dynamic slide")


if __name__ == "__main__":

    parser = argparse.ArgumentParser()
    parser.add_argument("--no-copy", action="store_true", help="Deactivate the run_copy function")
    args = parser.parse_args()

    # print(not args.no_copy)

    # finish_up(not args.no_copy)
    if not args.no_copy:
        finish_up()
    else:
        fix_svelte_routes()

